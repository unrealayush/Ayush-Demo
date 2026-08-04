#!/usr/bin/env python3
"""
generate_editable_deck.py — Generates a 1-Page Editable PowerPoint (.pptx) Deck
containing the full Mevreon Bio-AI In-Silico Docking Pipeline Process, Models, Logic,
and Editable Flowchart.
"""

import sys
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

BASE_DIR = Path(__file__).resolve().parent.parent

def create_presentation():
    prs = Presentation()
    # Set 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # ── Colors ──
    COLOR_BG = RGBColor(4, 7, 13)            # #04070D
    COLOR_CARD = RGBColor(12, 19, 32)        # #0C1320
    COLOR_CARD_BORDER = RGBColor(30, 58, 86)  # #1E3A56
    COLOR_CYAN = RGBColor(6, 182, 212)       # #06B6D4
    COLOR_TEAL = RGBColor(20, 184, 166)      # #14B8A6
    COLOR_BLUE = RGBColor(59, 130, 246)      # #3B82F6
    COLOR_PURPLE = RGBColor(168, 85, 247)    # #A855F7
    COLOR_TEXT_MAIN = RGBColor(241, 245, 249)# #F1F5F9
    COLOR_TEXT_SUB = RGBColor(148, 163, 184) # #94A3B8

    # ── Background ──
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_BG
    bg.line.fill.background()

    # ── Top Title Bar ──
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.133), Inches(0.9))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = "MEVREON BIO-AI IN-SILICO DOCKING PIPELINE"
    p.font.name = "Arial"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    p2 = tf.add_paragraph()
    p2.text = "End-to-End Molecular Docking, Generative AI Pose Confidence, Interaction Fingerprints & Evidence Passport Architecture"
    p2.font.name = "Arial"
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_TEXT_SUB
    p2.space_before = Pt(4)

    # ── Flowchart Stages Container ──
    stages = [
        {
            "num": "01",
            "title": "Ligand & Target Prep",
            "subtitle": "Structure Conversion & Cleaning",
            "color": COLOR_CYAN,
            "bullets": [
                "SMILES → 3D Conformers (RDKit ETKDGv3)",
                "PDBQT Format Conversion (OpenBabel)",
                "Receptor Cleaning & Centroid Grid Box"
            ]
        },
        {
            "num": "02",
            "title": "Dual AI-Physics Docking",
            "subtitle": "Thermodynamics & Generative AI",
            "color": COLOR_TEAL,
            "bullets": [
                "AutoDock Vina: Physics ΔG (kcal/mol)",
                "DiffDock-L: Generative Diffusion on GPU",
                "Pose Scoring & RMSD Cluster Analysis"
            ]
        },
        {
            "num": "03",
            "title": "Interactions & MoA Graph",
            "subtitle": "2D/3D Fingerprints & Pathways",
            "color": COLOR_BLUE,
            "bullets": [
                "3.5Å Non-Covalent Interaction Parser",
                "H-Bonds, Hydrophobic, Salt Bridges & π-stacking",
                "Mechanism of Action Cascade Graph Builder"
            ]
        },
        {
            "num": "04",
            "title": "Scoring & Evidence Passport",
            "subtitle": "Validation Score & 11 Deliverables",
            "color": COLOR_PURPLE,
            "bullets": [
                "Validation Priority Score (0-100)",
                "Preclinical Evidence Passport (.md & .json)",
                "11 Complete Deliverables Package Output"
            ]
        }
    ]

    card_width = Inches(2.75)
    card_height = Inches(2.4)
    start_x = Inches(0.6)
    gap_x = Inches(0.35)
    start_y = Inches(1.5)

    for i, stage in enumerate(stages):
        x_pos = start_x + i * (card_width + gap_x)

        # Card Box Shape
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, start_y, card_width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = stage["color"]
        card.line.width = Pt(1.5)

        # Header Text Box
        c_tf = card.text_frame
        c_tf.word_wrap = True
        c_tf.vertical_anchor = MSO_ANCHOR.TOP
        c_tf.margin_left = Inches(0.15)
        c_tf.margin_right = Inches(0.15)
        c_tf.margin_top = Inches(0.15)
        c_tf.margin_bottom = Inches(0.15)

        cp1 = c_tf.paragraphs[0]
        cp1.text = f"STAGE {stage['num']} — {stage['title'].upper()}"
        cp1.font.name = "Arial"
        cp1.font.size = Pt(10)
        cp1.font.bold = True
        cp1.font.color.rgb = stage["color"]

        cp2 = c_tf.add_paragraph()
        cp2.text = stage["subtitle"]
        cp2.font.name = "Arial"
        cp2.font.size = Pt(8.5)
        cp2.font.color.rgb = COLOR_TEXT_SUB
        cp2.space_before = Pt(2)
        cp2.space_after = Pt(8)

        for b in stage["bullets"]:
            bp = c_tf.add_paragraph()
            bp.text = f"• {b}"
            bp.font.name = "Arial"
            bp.font.size = Pt(9)
            bp.font.color.rgb = COLOR_TEXT_MAIN
            bp.space_before = Pt(3)

        # Connector Arrow (between stages)
        if i < len(stages) - 1:
            arrow_x = x_pos + card_width + Inches(0.05)
            arrow_y = start_y + card_height / 2 - Inches(0.15)
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, arrow_x, arrow_y, Inches(0.25), Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = stage["color"]
            arrow.line.fill.background()

    # ── Section 2: Detailed Architecture & Subsystems Breakdown ──
    sub_y = Inches(4.15)
    sub_height = Inches(2.2)

    subsystems = [
        {
            "title": "INPUT ENGINE",
            "desc": "Ligand Libraries & Receptors",
            "color": COLOR_CYAN,
            "details": "• 24 AYUSH Phytochemicals\n• 12 AMR Pathogen Targets\n• SMILES Canonical Strings\n• PDBQT Coordinate Boxes"
        },
        {
            "title": "DOCKING MODELS",
            "desc": "Thermodynamics & Generative AI",
            "color": COLOR_TEAL,
            "details": "• AutoDock Vina Physics Engine\n• DiffDock-L GPU Inference\n• Binding Energy (ΔG kcal/mol)\n• Generative Confidence (+0.792)"
        },
        {
            "title": "MECHANISM PARSER",
            "desc": "Fingerprints & MoA Graph",
            "color": COLOR_BLUE,
            "details": "• 3.5Å Contact Radius Engine\n• 2D/3D Interaction Types\n• Pathway Cascade Graphing\n• Phenotype Disruption Track"
        },
        {
            "title": "EVIDENCE ENGINE",
            "desc": "Prioritization & Passport",
            "color": COLOR_PURPLE,
            "details": "• Multi-Factor Scoring (0-100)\n• High Priority Lead Decision\n• 11 Deliverable Files Output\n• Preclinical Audit Passport"
        }
    ]

    for j, sub in enumerate(subsystems):
        sx_pos = start_x + j * (card_width + gap_x)

        sub_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, sx_pos, sub_y, card_width, sub_height
        )
        sub_card.fill.solid()
        sub_card.fill.fore_color.rgb = COLOR_CARD
        sub_card.line.color.rgb = COLOR_CARD_BORDER
        sub_card.line.width = Pt(1)

        s_tf = sub_card.text_frame
        s_tf.word_wrap = True
        s_tf.margin_left = Inches(0.15)
        s_tf.margin_right = Inches(0.15)
        s_tf.margin_top = Inches(0.15)
        s_tf.margin_bottom = Inches(0.15)

        sp1 = s_tf.paragraphs[0]
        sp1.text = sub["title"]
        sp1.font.name = "Arial"
        sp1.font.size = Pt(11)
        sp1.font.bold = True
        sp1.font.color.rgb = sub["color"]

        sp2 = s_tf.add_paragraph()
        sp2.text = sub["desc"]
        sp2.font.name = "Arial"
        sp2.font.size = Pt(8.5)
        sp2.font.color.rgb = COLOR_TEXT_SUB
        sp2.space_before = Pt(2)
        sp2.space_after = Pt(6)

        sp3 = s_tf.add_paragraph()
        sp3.text = sub["details"]
        sp3.font.name = "Arial"
        sp3.font.size = Pt(8.5)
        sp3.font.color.rgb = COLOR_TEXT_MAIN
        sp3.line_spacing = 1.2

    # ── Bottom Summary Status Bar ──
    bar_y = Inches(6.5)
    bar_height = Inches(0.6)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), bar_y, Inches(12.133), bar_height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(15, 23, 42)
    bar.line.color.rgb = COLOR_CYAN
    bar.line.width = Pt(1)

    b_tf = bar.text_frame
    b_tf.word_wrap = True
    b_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    bp = b_tf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    bp.text = "⚡ 288 DOCKING PAIRS PRE-COMPUTED  |  12 AMR TARGETS  |  24 AYUSH COMPOUNDS  |  11 FORENSIC DELIVERABLES  |  NVIDIA L4 GPU ACCELERATED"
    bp.font.name = "Arial"
    bp.font.size = Pt(9.5)
    bp.font.bold = True
    bp.font.color.rgb = COLOR_CYAN

    # Save outputs
    docs_dir = BASE_DIR / "docs"
    outputs_dir = BASE_DIR / "outputs"
    os.makedirs(docs_dir, exist_ok=True)
    os.makedirs(outputs_dir, exist_ok=True)

    file1 = docs_dir / "Mevreon_BioAI_Pipeline_Flowchart_Deck.pptx"
    file2 = outputs_dir / "Mevreon_BioAI_Pipeline_Flowchart_Deck.pptx"

    prs.save(str(file1))
    prs.save(str(file2))
    print(f"[SUCCESS] Presentation Deck saved to:\n  - {file1}\n  - {file2}")

if __name__ == "__main__":
    create_presentation()
